import argparse
import os
from typing import List, Dict
from flask import Flask, request, jsonify, g
import requests
import google.generativeai as genai
from sentence_transformers import SentenceTransformer
import pdfplumber
from pptx import Presentation
from PIL import Image
import pytesseract
import platform
from flask_cors import CORS
from transformers import T5ForConditionalGeneration, T5Tokenizer
import warnings
from pinecone import Pinecone
from pinecone import ServerlessSpec
import time

# Suppress specific warnings from Hugging Face transformers library
warnings.filterwarnings("ignore", message="You are using the default legacy behaviour")
warnings.filterwarnings("ignore", message="It will be set to `False` by default.")
warnings.filterwarnings("ignore", message="`clean_up_tokenization_spaces` was not set")

# Set Google API key for Gemini
def set_google_api_key():
    api_key = "AIzaSyD7VrRJrSa3W7u0syiZpWldChRCTiWLp-4"
    if platform.system() == "Windows":
        os.environ["GOOGLE_API_KEY"] = api_key
    else:  # Assuming macOS or Linux
        os.environ["GOOGLE_API_KEY"] = api_key

set_google_api_key()

# Initialize Flask app
app = Flask(__name__)
CORS(app)
os.environ["TOKENIZERS_PARALLELISM"] = "false"
# Retrieve the API key from the environment variable when needed
google_api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=google_api_key)

# Load the Flan-T5 model and tokenizer from Hugging Face
flan_tokenizer = T5Tokenizer.from_pretrained("google/flan-t5-base", legacy=False)
flan_model = T5ForConditionalGeneration.from_pretrained("google/flan-t5-base")

# Initialize history storage
session_history: Dict[str, List[Dict[str, str]]] = {}

# Load Huggingface model for embeddings
embedding_model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

api_key = "7968bf7e-97f5-4022-adf7-9294590606be"

# Configure Pinecone client
pc = Pinecone(api_key=api_key)
cloud = os.environ.get('PINECONE_CLOUD') or 'aws'
region = os.environ.get('PINECONE_REGION') or 'us-east-1'
spec = ServerlessSpec(cloud=cloud, region=region)
index_name = 'msnabiel'
existing_indexes = [
    index_info["name"] for index_info in pc.list_indexes()
]

# Check if index already exists (it shouldn't if this is the first time)
if index_name not in existing_indexes:
    # If does not exist, create index
    pc.create_index(
        index_name,
        dimension=384,  # Dimensionality of minilm
        metric='cosine',  # Use cosine similarity
        spec=spec
    )
    # Wait for index to be initialized
    while not pc.describe_index(index_name).status['ready']:
        time.sleep(1)

# Connect to index
index = pc.Index(index_name)
time.sleep(1)

# Initialize Gemini model
model = genai.GenerativeModel("gemini-1.5-flash")

# Helper functions to extract text from different document types
def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file."""
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_image(image_path):
    """Extract text from an image using OCR."""
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    return text

def extract_text_from_file(file_path):
    """Extract text based on the file type."""
    _, extension = os.path.splitext(file_path)
    if extension.lower() == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension.lower() in [".ppt", ".pptx"]:
        return extract_text_from_ppt(file_path)
    elif extension.lower() in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
        return extract_text_from_image(file_path)
    else:
        return None

def fetch_and_call_api(query):
    response_data = {}
    try:
        response = requests.get("https://dummyapi.com/api", params={"query": query})
        response_data = response.json()
    except Exception as e:
        return "Need API Key to call, to perform the action. "

    if response and response_data.get("status") == "success":
        return response_data["message"]

def execute_action(action_name: str) -> str:
    if action_name == "create_order":
        fetch_response = fetch_and_call_api("YOUR_API_KEY")
        return fetch_response + " Order created successfully."
    elif action_name == "cancel_order":
        fetch_response = fetch_and_call_api("YOUR_API_KEY")
        return fetch_response + " Order cancelled successfully."
    elif action_name == "collect_payment":
        fetch_response = fetch_and_call_api("YOUR_API_KEY")
        return fetch_response + " Payment collected successfully."
    elif action_name == "view_invoice":
        fetch_response = fetch_and_call_api("YOUR_API_KEY")
        return fetch_response + " Here is your invoice."
    else:
        return "No action taken."

# Define the default actions list outside the function
DEFAULT_ACTIONS_LIST = ["create_order", "cancel_order", "collect_payment", "view_invoice"]

def classify_query_with_flan(query: str, actions_list: list = DEFAULT_ACTIONS_LIST) -> str:
    """Use Flan-T5 to classify the query as action-based or context-based using a dynamic list of actions."""
    
    # Convert the list of actions to a string to use in the prompt
    actions_str = ', '.join(actions_list)
    
    # Dynamic prompt using the actions list
    prompt = f"""
    Classify the following query strictly as one of the actions: {actions_str}, or context-based.
    
    Query: {query}
    """
    print("Query for FlanT5:", query)
    # Tokenize the input
    inputs = flan_tokenizer(prompt, return_tensors="pt")
    outputs = flan_model.generate(**inputs, max_length=10, num_return_sequences=1)
    response = flan_tokenizer.decode(outputs[0], skip_special_tokens=True).lower()
    print("Response from FlanT5:", response)
    
    # Check if the response matches one of the actions
    for action in actions_list:
        if action in response:
            return action
    return "context_based"

def build_combined_prompt(query: str, context: List[str], history: List[Dict[str, str]]) -> str:
    base_prompt = """
        I am going to ask you a question, which I would like you to answer strictly based on the given context.
        If there is not enough information in the context to answer the question, make a guess based on the context.
        """
    user_prompt = f"The question is '{query}'. Here is all the context you have: {' '.join(context)}"
    history_prompt = "\n".join([f"User: {item['query']}\nBot: {item['response']}" for item in history])

    print("Combined Prompt:", f"{base_prompt} {history_prompt} {user_prompt}")

    return f"{base_prompt} {history_prompt} {user_prompt}"

def get_gemini_response(query: str, context: List[str], session_id: str, document_id: str) -> str:
    history = session_history.get(session_id, [])

    # Classify the query using Flan-T5
    action = classify_query_with_flan(query)
    action = action.lower()
    if action != "context_based":
        action_response = execute_action(action)
        
        # Append the action execution to session history
        session_history.setdefault(session_id, []).append({"query": query, "response": action_response})
        return action_response

    # Build the combined prompt
    prompt = build_combined_prompt(query, context, history)

    # Get response from Gemini
    response = model.generate_content(prompt)

    response_text = response.text.strip().lower()

    # Save the query and response in session history
    session_history.setdefault(session_id, []).append({"query": query, "response": response.text})

    # Add references to the response
    references = "\n".join([f"From document '{document_id}': Line {i + 1}: {line}" 
                            for i, line in enumerate(context)])

    return f"{response.text}\n\nReferences:\n{references}"

@app.route('/upload_document', methods=['POST'])
def upload_document():
    # Check if the uploads directory exists, if not create it
    uploads_dir = "uploads"
    os.makedirs(uploads_dir, exist_ok=True)  # This will create the directory if it doesn't exist

    if 'document' not in request.files:
        return jsonify({"error": "No document uploaded."}), 400

    document = request.files['document']
    doc_name = document.filename

    # Validate document format
    if not doc_name.endswith(('.pdf', '.ppt', '.pptx', '.jpg', '.jpeg', '.png', '.bmp', '.tiff')):
        return jsonify({"error": "Unsupported file format."}), 400

    # Save the document
    file_path = os.path.join(uploads_dir, doc_name)
    document.save(file_path)

    # Extract text from the document
    extracted_text = extract_text_from_file(file_path)

    if extracted_text is None:
        return jsonify({"error": "Unable to extract text from the document."}), 400

    # Create vectors and store in Pinecone
    vectors = embedding_model.encode(extracted_text.splitlines()).tolist()
    for i, line in enumerate(extracted_text.splitlines()):
        index.upsert([(f"{doc_name}-{i}", vectors[i], {"line": line})])

    return jsonify({"message": "Document uploaded and processed successfully."}), 200

@app.route('/chat', methods=['POST'])
def query():
    data = request.get_json()
    query_text = data.get("query")
    session_id = data.get("session_id")
    document_id = data.get("document_id")  # New parameter

    if not query_text or not session_id:
        return jsonify({"error": "Missing query or session ID."}), 400

    # Retrieve the Pinecone vectors to provide context
    query_vector = embedding_model.encode(query_text).tolist()

    # Adjust query based on document_id
    if document_id:  # If document_id is provided
        # Query for lines associated with the specific document
        context_vectors = index.query(vector=query_vector, top_k=5, include_metadata=True)
        
        # Filter context to include only lines from the specified document ID
        context = [match["metadata"].get("line", "No text available") 
                   for match in context_vectors["matches"] 
                   if match["id"].startswith(document_id)]
        if not context:
            return jsonify({"error": "No relevant context found for the specified document."}), 404
    else:  # No document_id provided
        # Retrieve context from the entire index
        context_vectors = index.query(vector=query_vector, top_k=5, include_metadata=True)
        # Extract context lines without filtering by document ID
        context = [match["metadata"].get("line", "No text available") for match in context_vectors["matches"]]

    print("Context:", context)

    # Generate response using Gemini with document_id
    response = get_gemini_response(query_text, context, session_id, document_id)

    return jsonify({"response": response}), 200

if __name__ == '__main__':
    app.run(debug=True, port=5000)
